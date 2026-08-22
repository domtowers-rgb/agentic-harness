import re

from pptx import Presentation

from plugins.file_ops import SANDBOX_DIR, _resolve_safe

MAX_SLIDES = 100
MAX_BULLETS_PER_SLIDE = 50


def _sanitize_filename(name: str) -> str:
    name = re.sub(r"[^A-Za-z0-9 _-]", "", name).strip().replace(" ", "_")
    name = name[:80] or "presentation"
    if not name.lower().endswith(".pptx"):
        name += ".pptx"
    return name


def create_presentation(title: str, slides: list, subtitle: str = None, filename: str = None):
    """Create a .pptx file in the sandboxed working directory: a title slide
    followed by one title+bullets slide per entry in `slides`."""
    if not title or not isinstance(title, str):
        return {"error": "title is required"}
    if not isinstance(slides, list):
        return {"error": "slides must be a list of {title, bullets} objects"}
    if len(slides) > MAX_SLIDES:
        return {"error": f"too many slides (max {MAX_SLIDES})"}

    try:
        target = _resolve_safe(_sanitize_filename(filename or title))
    except ValueError as exc:
        return {"error": str(exc)}

    try:
        prs = Presentation()

        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = title
        if subtitle and len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = subtitle

        bullet_layout = prs.slide_layouts[1]
        for i, entry in enumerate(slides):
            if not isinstance(entry, dict):
                return {"error": f"slide {i} must be an object with 'title' and optional 'bullets'"}
            bullets = entry.get("bullets") or []
            if len(bullets) > MAX_BULLETS_PER_SLIDE:
                return {"error": f"slide {i} has too many bullets (max {MAX_BULLETS_PER_SLIDE})"}

            slide = prs.slides.add_slide(bullet_layout)
            slide.shapes.title.text = str(entry.get("title") or "")
            if bullets and len(slide.placeholders) > 1:
                text_frame = slide.placeholders[1].text_frame
                text_frame.text = str(bullets[0])
                for bullet in bullets[1:]:
                    text_frame.add_paragraph().text = str(bullet)

        target.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(target))
    except Exception as exc:
        return {"error": f"failed to create presentation: {exc}"}

    return {
        "status": "written",
        "path": str(target.relative_to(SANDBOX_DIR)),
        "slide_count": len(slides) + 1,
    }


def register(registry):
    registry.register("create_presentation", create_presentation, {
        "name": "create_presentation",
        "description": (
            f"Create a PowerPoint (.pptx) presentation in the sandboxed working directory ({SANDBOX_DIR}). "
            "The first slide is a title slide; each entry in 'slides' becomes a slide with a title and "
            "bullet points."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "title": {"type": "string", "description": "Presentation title, shown on the first slide"},
                "subtitle": {"type": "string", "description": "Optional subtitle shown under the title on the first slide"},
                "slides": {
                    "type": "array",
                    "description": "Content slides after the title slide",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string"},
                            "bullets": {"type": "array", "items": {"type": "string"}},
                        },
                        "required": ["title"],
                    },
                },
                "filename": {
                    "type": "string",
                    "description": "Output filename, e.g. 'my_deck.pptx'. Defaults to a name derived from the title.",
                },
            },
            "required": ["title", "slides"],
        },
    })
